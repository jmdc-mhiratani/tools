"""
Conversion implementations using the refactored core modules.
Handles image and PPTX conversion with proper error handling and progress tracking.
"""

from __future__ import annotations

from pathlib import Path
from typing import List, Optional
from io import BytesIO

from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE

from ..core.pdf_processor import (
    ConversionConfig,
    open_pdf_document,
    process_page_to_pixmap,
    process_page_to_bytes,
    count_total_pages,
    mm_to_emu,
    points_to_emu
)
from ..utils.error_handling import handle_pdf_errors, handle_file_errors, UserFriendlyError
from ..utils.path_utils import PathManager


class BaseConverter:
    """
    Base class for PDF converters with common functionality.
    """

    def __init__(self, path_manager: PathManager, progress_tracker):
        self.path_manager = path_manager
        self.progress_tracker = progress_tracker

    def _get_pdf_files(self) -> List[Path]:
        """
        Get list of PDF files to process.

        Returns:
            List of PDF file paths

        Raises:
            UserFriendlyError: If no PDF files found
        """
        pdf_files = self.path_manager.find_pdf_files()
        if not pdf_files:
            raise UserFriendlyError(
                message="Inputフォルダに変換するPDFファイルが見つかりません",
                suggestion="PDFファイルをInputフォルダに配置してください"
            )
        return pdf_files

    def _setup_progress(self, pdf_files: List[Path]) -> None:
        """Set up progress tracking for conversion."""
        total_pages = count_total_pages(pdf_files)
        self.progress_tracker.set_maximum(total_pages)


class ImageConverter(BaseConverter):
    """
    Converts PDF files to PNG images.
    """

    @handle_pdf_errors
    @handle_file_errors
    def convert_all_pdfs(self, config: ConversionConfig) -> List[Path]:
        """
        Convert all PDFs in input directory to PNG images.

        Args:
            config: Conversion configuration

        Returns:
            List of generated image file paths

        Raises:
            UserFriendlyError: If conversion fails
        """
        pdf_files = self._get_pdf_files()
        self._setup_progress(pdf_files)

        output_files = []

        for pdf_file in pdf_files:
            pdf_output_files = self._convert_single_pdf(pdf_file, config)
            output_files.extend(pdf_output_files)

            # Check for cancellation
            if self.progress_tracker.is_cancelled:
                break

        return output_files

    def _convert_single_pdf(self, pdf_file: Path, config: ConversionConfig) -> List[Path]:
        """
        Convert a single PDF file to PNG images.

        Args:
            pdf_file: PDF file to convert
            config: Conversion configuration

        Returns:
            List of generated image file paths
        """
        output_files = []
        base_name = pdf_file.stem

        with open_pdf_document(pdf_file) as doc:
            for page_num, page in enumerate(doc, start=1):
                # Process page to pixmap
                pixmap, page_info = process_page_to_pixmap(page, config)

                try:
                    # Generate output filename
                    output_filename = f"{base_name}_{page_num}"
                    output_path = self.path_manager.get_output_path(
                        output_filename, "png"
                    )

                    # Save image
                    pixmap.save(str(output_path))
                    output_files.append(output_path)

                    # Update progress
                    self.progress_tracker.step(1)

                finally:
                    # Clean up pixmap memory
                    pixmap = None

        return output_files


class PPTXConverter(BaseConverter):
    """
    Converts PDF files to a single PPTX presentation.
    """

    @handle_pdf_errors
    @handle_file_errors
    def convert_all_pdfs(self, config: ConversionConfig) -> Optional[Path]:
        """
        Convert all PDFs to a single PPTX presentation.

        Args:
            config: Conversion configuration

        Returns:
            Path to generated PPTX file, or None if no files processed

        Raises:
            UserFriendlyError: If conversion fails
        """
        pdf_files = self._get_pdf_files()
        self._setup_progress(pdf_files)

        # Create presentation
        presentation = self._create_presentation(config)

        # Process all PDF files
        first_pdf_name = pdf_files[0].stem
        slide_count = 0

        for pdf_file in pdf_files:
            slide_count += self._add_pdf_to_presentation(
                pdf_file, presentation, config
            )

            # Check for cancellation
            if self.progress_tracker.is_cancelled:
                break

        if slide_count == 0:
            return None

        # Save presentation
        output_path = self.path_manager.get_output_path(
            first_pdf_name, "pptx"
        )
        presentation.save(str(output_path))

        return output_path

    def _create_presentation(self, config: ConversionConfig) -> Presentation:
        """
        Create a new PowerPoint presentation with A3 landscape dimensions.

        Args:
            config: Conversion configuration

        Returns:
            Configured presentation object
        """
        presentation = Presentation()

        # Set slide dimensions to A3 landscape (420mm x 297mm)
        presentation.slide_width = mm_to_emu(config.slide_width_mm)
        presentation.slide_height = mm_to_emu(config.slide_height_mm)

        return presentation

    def _add_pdf_to_presentation(
        self,
        pdf_file: Path,
        presentation: Presentation,
        config: ConversionConfig
    ) -> int:
        """
        Add all pages from a PDF file to the presentation.

        Args:
            pdf_file: PDF file to add
            presentation: Target presentation
            config: Conversion configuration

        Returns:
            Number of slides added
        """
        slides_added = 0
        base_name = pdf_file.stem

        with open_pdf_document(pdf_file) as doc:
            for page_num, page in enumerate(doc, start=1):
                self._add_page_to_presentation(
                    page, presentation, config, base_name, page_num
                )
                slides_added += 1
                self.progress_tracker.step(1)

        return slides_added

    def _add_page_to_presentation(
        self,
        page,
        presentation: Presentation,
        config: ConversionConfig,
        base_name: str,
        page_num: int
    ) -> None:
        """
        Add a single PDF page to the presentation.

        Args:
            page: PyMuPDF page object
            presentation: Target presentation
            config: Conversion configuration
            base_name: Base filename for labeling
            page_num: Page number for labeling
        """
        # Get page dimensions and determine rotation
        rect = page.rect
        is_portrait = rect.width < rect.height
        rotation = 90 if (config.auto_rotate and is_portrait) else 0

        # Calculate image dimensions after rotation
        if rotation == 90:
            width_emu = points_to_emu(rect.height)
            height_emu = points_to_emu(rect.width)
        else:
            width_emu = points_to_emu(rect.width)
            height_emu = points_to_emu(rect.height)

        # Generate high-quality image
        image_bytes, _ = process_page_to_bytes(page, config, "png")
        image_stream = BytesIO(image_bytes)

        # Add slide
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])  # Blank layout

        # Calculate positioning to center image on slide
        slide_width = presentation.slide_width
        slide_height = presentation.slide_height

        final_width, final_height = self._calculate_fitted_dimensions(
            width_emu, height_emu, slide_width, slide_height
        )

        left = int((slide_width - final_width) / 2)
        top = int((slide_height - final_height) / 2)

        # Add image to slide
        slide.shapes.add_picture(
            image_stream, left, top, width=final_width, height=final_height
        )

        # Add filename label
        self._add_filename_label(slide, f"{base_name}_{page_num}", presentation)

    def _calculate_fitted_dimensions(
        self,
        image_width: int,
        image_height: int,
        slide_width: int,
        slide_height: int
    ) -> tuple[int, int]:
        """
        Calculate dimensions to fit image within slide while maintaining aspect ratio.

        Args:
            image_width: Original image width in EMU
            image_height: Original image height in EMU
            slide_width: Slide width in EMU
            slide_height: Slide height in EMU

        Returns:
            Tuple of (fitted_width, fitted_height) in EMU
        """
        # If image fits within slide, use original size
        if image_width <= slide_width and image_height <= slide_height:
            return image_width, image_height

        # Calculate scaling ratios
        width_ratio = slide_width / image_width
        height_ratio = slide_height / image_height

        # Use the smaller ratio to ensure image fits
        scale_ratio = min(width_ratio, height_ratio)

        fitted_width = int(image_width * scale_ratio)
        fitted_height = int(image_height * scale_ratio)

        return fitted_width, fitted_height

    def _add_filename_label(
        self,
        slide,
        label_text: str,
        presentation: Presentation
    ) -> None:
        """
        Add filename label to the top-left corner of the slide.

        Args:
            slide: Slide to add label to
            label_text: Text for the label
            presentation: Presentation for dimensions
        """
        # Label dimensions (30% of slide width, 6% of slide height)
        label_width = int(presentation.slide_width * 0.3)
        label_height = int(presentation.slide_height * 0.06)

        # Create text box
        textbox = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            0, 0,  # Top-left corner
            label_width,
            label_height
        )

        # Style the text box
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 102, 0)  # Orange background
        textbox.line.color.rgb = RGBColor(255, 0, 0)  # Red border
        textbox.shadow.inherit = False

        # Set text properties
        text_frame = textbox.text_frame
        paragraph = text_frame.paragraphs[0]
        paragraph.text = label_text
        paragraph.font.name = "游ゴシック"
        paragraph.font.size = Pt(18)
        paragraph.font.bold = True
        paragraph.font.color.rgb = RGBColor(255, 255, 255)  # White text
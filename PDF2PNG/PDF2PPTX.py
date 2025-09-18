"""
PDF2PPTX.py
============================================
PDF の各ページを "A3 横 (420×297 mm)" スライドに実寸で貼り付ける GUI ツール。
縦長でスキャンされた PDF ページは自動で 90° 回転して横向きに配置します。

主要機能
---------
1. **PDF → PNG 変換**
   *Input* フォルダの PDF を PNG へ一括変換（スケール指定可）。
   ポートレート判定時は横向き画像へ自動回転。

2. **PDF → PPTX 直接変換**
   *Input* の全 PDF を読み込み、各ページを A3 横スライド中央へ実寸配置。
   ポートレート判定なら 90° 回転し横向きに挿入。

3. **Input/Output フォルダ初期化**

動作環境
---------
* Python 3.10+
* PyMuPDF, python-pptx, Pillow

"""

from __future__ import annotations

import sys
import os
import shutil
import webbrowser
from io import BytesIO
from typing import Tuple

import fitz  # PyMuPDF
import tkinter as tk
from tkinter import messagebox, ttk
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image

# ---------------------------------------------
# 設定
# ---------------------------------------------
# 縦長ページを自動で 90° 回転して横向きで貼り付ける
AUTO_ROTATE_TO_LANDSCAPE = True

# ---------------------------------------------
# ヘルパー関数
# ---------------------------------------------

def mm_to_emu(mm: float) -> int:
    """ミリメートル → EMU"""
    return int((mm / 25.4) * 914400)


def points_to_emu(pt: float) -> int:
    """PDF の pt → EMU"""
    return int((pt / 72) * 914400)


def get_app_path() -> str:
    """exe / スクリプト配置ディレクトリを返す (PyInstaller 対応)"""
    if hasattr(sys, "_MEIPASS"):
        return os.path.dirname(sys.executable)
    return os.path.dirname(os.path.abspath(__file__))


def ensure_dirs() -> None:
    base = get_app_path()
    for name in ("Input", "Output"):
        d = os.path.join(base, name)
        os.makedirs(d, exist_ok=True)


# ---------------------------------------------
# 機能 1：PDF → PNG
# ---------------------------------------------

def rotate_if_portrait(page: fitz.Page, pix: fitz.Pixmap) -> Tuple[fitz.Pixmap, bool]:
    """ポートレートなら横向きに 90° 回転。戻り値: (Pixmap, rotated?)"""
    portrait = page.rect.width < page.rect.height
    if AUTO_ROTATE_TO_LANDSCAPE and portrait:
        return page.get_pixmap(matrix=fitz.Matrix(scale_factor, scale_factor), rotate=90), True
    return pix, False


def convert_pdfs_to_images() -> None:
    base = get_app_path()
    input_dir = os.path.join(base, "Input")
    output_dir = os.path.join(base, "Output")
    os.makedirs(output_dir, exist_ok=True)

    pdf_files = [f for f in os.listdir(input_dir) if f.lower().endswith(".pdf")]
    if not pdf_files:
        messagebox.showwarning("警告", "Input フォルダに PDF が見つかりません")
        return

    # スケール倍率取得
    global scale_factor
    try:
        scale_factor = float(entry_scale.get())
    except ValueError:
        messagebox.showerror("エラー", "スケールは数値で指定してください")
        return

    # ページ総数
    total_pages = sum(len(fitz.open(os.path.join(input_dir, f))) for f in pdf_files)
    progress.configure(maximum=total_pages, value=0)
    root.update_idletasks()

    for pdf_name in pdf_files:
        doc = fitz.open(os.path.join(input_dir, pdf_name))
        base_name = os.path.splitext(pdf_name)[0]

        for i, page in enumerate(doc):
            # 初期 Pixmap
            pix = page.get_pixmap(matrix=fitz.Matrix(scale_factor, scale_factor))
            pix, rotated = rotate_if_portrait(page, pix)

            out_name = f"{base_name}_{i+1}.png"
            pix.save(os.path.join(output_dir, out_name))

            progress.step(1)
            root.update_idletasks()

    messagebox.showinfo("完了", f"PNG 変換完了\n\n保存先: {output_dir}")
    webbrowser.open(output_dir)
    progress.configure(value=0)


# ---------------------------------------------
# 機能 2：PDF → PPTX 直接変換
# ---------------------------------------------

def create_pptx_direct_from_pdfs() -> None:
    base = get_app_path()
    input_dir = os.path.join(base, "Input")
    pdf_files = [f for f in os.listdir(input_dir) if f.lower().endswith(".pdf")]
    if not pdf_files:
        messagebox.showwarning("警告", "Input フォルダに PDF が見つかりません")
        return

    try:
        scale = float(entry_scale.get())
    except ValueError:
        messagebox.showerror("エラー", "スケールは数値で指定してください")
        return

    prs = Presentation()
    prs.slide_width = mm_to_emu(420)
    prs.slide_height = mm_to_emu(297)

    total_pages = sum(len(fitz.open(os.path.join(input_dir, f))) for f in pdf_files)
    progress.configure(maximum=total_pages, value=0)
    root.update_idletasks()

    first_pdf_name = os.path.splitext(pdf_files[0])[0]

    for pdf_name in pdf_files:
        doc = fitz.open(os.path.join(input_dir, pdf_name))
        base_name = os.path.splitext(pdf_name)[0]

        for page_num, page in enumerate(doc, start=1):
            # 元ページ矩形
            rect = page.rect
            portrait = rect.width < rect.height
            rotate_deg = 90 if AUTO_ROTATE_TO_LANDSCAPE and portrait else 0

            # サイズ (回転後なら幅・高さを入れ替え)
            if rotate_deg == 90:
                width_emu = points_to_emu(rect.height)
                height_emu = points_to_emu(rect.width)
            else:
                width_emu = points_to_emu(rect.width)
                height_emu = points_to_emu(rect.height)

            # 画像化 (高解像度確保用 scale 倍率 + 回転)
            pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale), rotate=rotate_deg)
            img_bytes = pix.tobytes("png")
            img_stream = BytesIO(img_bytes)

            # スライド追加
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            sw, sh = prs.slide_width, prs.slide_height

            # スライドに収まるよう調整
            final_w, final_h = width_emu, height_emu
            if final_w > sw or final_h > sh:
                ratio = final_w / final_h
                slide_ratio = sw / sh
                if ratio > slide_ratio:
                    final_w = sw
                    final_h = int(sw / ratio)
                else:
                    final_h = sh
                    final_w = int(sh * ratio)
            left = int((sw - final_w) / 2)
            top = int((sh - final_h) / 2)

            slide.shapes.add_picture(img_stream, left, top, width=final_w, height=final_h)

            # 左上のファイル名ラベル
            label = f"{base_name}_{page_num}"
            tb = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, int(sw * 0.3), int(sh * 0.06)
            )
            tb.fill.solid()
            tb.fill.fore_color.rgb = RGBColor(255, 102, 0)
            tb.line.color.rgb = RGBColor(255, 0, 0)
            tb.shadow.inherit = False
            tf = tb.text_frame
            para = tf.paragraphs[0]
            para.text = label
            para.font.name = "游ゴシック"
            para.font.size = Pt(18)
            para.font.bold = True
            para.font.color.rgb = RGBColor(255, 255, 255)

            progress.step(1)
            root.update_idletasks()

    save_path = os.path.join(base, "Output", f"{first_pdf_name}.pptx")
    os.makedirs(os.path.dirname(save_path), exist_ok=True)
    prs.save(save_path)

    messagebox.showinfo("完了", f"PPTX 作成が完了しました\n\n保存先: {save_path}")
    webbrowser.open(save_path)
    progress.configure(value=0)


# ---------------------------------------------
# 機能 3：フォルダ初期化
# ---------------------------------------------

def clear_folders() -> None:
    base = get_app_path()
    for name in ("Input", "Output"):
        target = os.path.join(base, name)
        shutil.rmtree(target, ignore_errors=True)
        os.makedirs(target, exist_ok=True)
    messagebox.showinfo("初期化", "Input / Output フォルダを初期化しました")


# ---------------------------------------------
# GUI 構築
# ---------------------------------------------

ensure_dirs()

root = tk.Tk()
root.title("PDF2PPTX Converter")

# スケール入力
tk.Label(root, text="スケール倍率 (例: 1.5)").grid(row=0, column=0, sticky="e", padx=8, pady=8)
entry_scale = tk.Entry(root, width=8)
entry_scale.insert(0, "1.5")
entry_scale.grid(row=0, column=1, padx=8, pady=8, sticky="w")

# 自動回転チェック
auto_rotate_var = tk.BooleanVar(value=AUTO_ROTATE_TO_LANDSCAPE)

def toggle_auto_rotate() -> None:
    global AUTO_ROTATE_TO_LANDSCAPE
    AUTO_ROTATE_TO_LANDSCAPE = auto_rotate_var.get()

auto_chk = tk.Checkbutton(root, text="縦長ページを横向きに自動回転", variable=auto_rotate_var, command=toggle_auto_rotate)
auto_chk.grid(row=1, column=0, columnspan=2, pady=(0, 8))

# ボタン
btn_pdf2png = tk.Button(root, text="📄 PDF → PNG 変換", width=30, command=convert_pdfs_to_images)
btn_pdf2png.grid(row=2, column=0, columnspan=2, pady=4)

btn_pdf2ppt = tk.Button(root, text="📈 PDF → PPTX 変換 (A3 横)", width=30, command=create_pptx_direct_from_pdfs)
btn_pdf2ppt.grid(row=3, column=0, columnspan=2, pady=4)

btn_clear = tk.Button(root, text="🧹 Input / Output 初期化", width=30, command=clear_folders)
btn_clear.grid(row=4, column=0, columnspan=2, pady=4)

# プログレスバー
progress = ttk.Progressbar(root, orient="horizontal", length=380, mode="determinate")
progress.grid(row=5, column=0, columnspan=2, pady=10)

root.mainloop()

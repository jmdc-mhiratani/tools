import tkinter as tk
from tkinter import messagebox, ttk
import fitz  # PyMuPDF
import os
import shutil
import webbrowser
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from PIL import Image
from io import BytesIO

# ===============================================
# ヘルパー関数
# ===============================================

def mm_to_emu(mm: float) -> int:
    """
    ミリメートルを EMU (PowerPoint の内部単位) に変換
    1 inch = 914400 EMU, 1 inch = 25.4 mm
    """
    return int((mm / 25.4) * 914400)

def points_to_emu(pt: float) -> int:
    """
    PDF のポイント(pt)を EMU に変換
    1 point = 1/72 inch, 1 inch = 914400 EMU
    """
    return int((pt / 72) * 914400)

def get_app_path() -> str:
    """ このスクリプトが置かれているディレクトリの絶対パス """
    return os.path.dirname(os.path.abspath(__file__))

def ensure_dirs():
    """
    Input, Output フォルダをアプリケーションフォルダ直下に作成(または存在チェック)。
    """
    base = get_app_path()
    for name in ("Input", "Output"):
        d = os.path.join(base, name)
        if not os.path.exists(d):
            os.makedirs(d)

# ===============================================
# 機能 1：PDF→PNG（一括スケール変換）
# ===============================================

def convert_pdfs_to_images():
    """
    Input フォルダ内のすべてのPDFを、指定スケールでPNGに変換し、Outputフォルダに保存する。
    プログレスバーでページ数ごとの進捗を表示。
    """
    base_dir = get_app_path()
    input_dir = os.path.join(base_dir, "Input")
    output_dir = os.path.join(base_dir, "Output")
    os.makedirs(output_dir, exist_ok=True)

    # PDFファイル一覧取得
    pdf_files = [f for f in os.listdir(input_dir) if f.lower().endswith(".pdf")]
    if not pdf_files:
        messagebox.showwarning("警告", "InputフォルダにPDFが見つかりません")
        return

    # 全ページ数をカウント
    total_pages = 0
    for fname in pdf_files:
        try:
            doc = fitz.open(os.path.join(input_dir, fname))
            total_pages += len(doc)
        except:
            continue

    # スケール取得
    try:
        scale = float(entry_scale.get())
    except ValueError:
        messagebox.showerror("エラー", "スケールは数値で指定してください")
        return

    # プログレスバー初期化
    progress["value"] = 0
    progress["maximum"] = total_pages
    root.update_idletasks()

    # 各PDF→PNG
    for file_name in pdf_files:
        file_path = os.path.join(input_dir, file_name)
        try:
            doc = fitz.open(file_path)
            base_name = os.path.splitext(file_name)[0]
            for page_number in range(len(doc)):
                page = doc[page_number]
                zoom_matrix = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=zoom_matrix)

                # 保存先パス
                output_file = os.path.join(output_dir, f"{base_name}_{page_number + 1}.png")
                pix.save(output_file)

                # プログレス更新
                progress["value"] += 1
                root.update_idletasks()
        except Exception as e:
            print(f"Error processing {file_name}: {e}")

    messagebox.showinfo("完了", f"画像変換が完了しました\n\n保存先: {output_dir}")
    webbrowser.open(output_dir)
    progress["value"] = 0

# ===============================================
# 機能 2：PNG→PPTX (A3 横スライド、元PDF実寸)
# ===============================================

def create_pptx_from_images():
    """
    Outputフォルダ内のPNG/JPG を、対応するPDFページの実寸サイズで
    A3横スライドに中央配置＋拡張子を除いたファイル名表示(左上)で
    1ファイルのPPTXにまとめる。
    """
    base_dir = get_app_path()
    input_dir = os.path.join(base_dir, "Input")
    output_dir = os.path.join(base_dir, "Output")
    image_files = sorted(
        f for f in os.listdir(output_dir)
        if f.lower().endswith((".png", ".jpg", ".jpeg"))
    )
    if not image_files:
        messagebox.showwarning("警告", "Outputフォルダに画像が見つかりません")
        return

    # プレゼンテーション作成＆A3 横にスライドサイズ設定
    prs = Presentation()
    prs.slide_width  = mm_to_emu(420)   # 420mm
    prs.slide_height = mm_to_emu(297)   # 297mm

    for image_name in image_files:
        # 画像ファイルのパス
        image_path = os.path.join(output_dir, image_name)

        # 対応するPDFファイルとページ番号を推測
        # 例: foo_3.png → foo.pdf の 3ページ目
        base_part, page_part = image_name.rsplit("_", 1)
        page_index_str = page_part.rsplit(".", 1)[0]  # "3"
        try:
            page_index = int(page_index_str) - 1  # 0-based index
        except:
            page_index = None

        pdf_path = os.path.join(input_dir, f"{base_part}.pdf")
        if page_index is not None and os.path.isfile(pdf_path):
            try:
                doc = fitz.open(pdf_path)
                page = doc[page_index]
                # PDFページの実寸サイズ（ポイント単位）
                rect = page.rect
                width_pt  = rect.width
                height_pt = rect.height
                # EMU に変換
                width_emu  = points_to_emu(width_pt)
                height_emu = points_to_emu(height_pt)
            except:
                width_emu  = None
                height_emu = None
        else:
            width_emu  = None
            height_emu = None

        # スライド追加
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank

        # 画像をメモリ経由で読み込み（BytesIO）すると、後で物理サイズ指定ができる
        with Image.open(image_path) as pil_img:
            img_byte_arr = BytesIO()
            pil_img.save(img_byte_arr, format='PNG')
            img_byte_arr.seek(0)

            # もし実寸サイズを取れていれば、そのサイズで配置。そうでなければ全画面フィット
            sw = prs.slide_width
            sh = prs.slide_height
            if width_emu and height_emu:
                final_w = width_emu
                final_h = height_emu
                # ページがスライドより大きい場合は縮小してフィット
                if final_w > sw or final_h > sh:
                    ratio = final_w / final_h
                    slide_ratio = sw / sh
                    if ratio > slide_ratio:
                        # 幅優先
                        final_w = sw
                        final_h = int(sw / ratio)
                    else:
                        final_h = sh
                        final_w = int(sh * ratio)
                left = int((sw - final_w) / 2)
                top  = int((sh - final_h) / 2)
            else:
                # 実寸取得できない場合は、画像をスライド最大にフィット
                iw, ih = pil_img.size
                img_w_emu = int((iw / 96) * 914400)
                img_h_emu = int((ih / 96) * 914400)
                ratio = img_w_emu / img_h_emu
                slide_ratio = sw / sh
                if ratio > slide_ratio:
                    final_w = sw
                    final_h = int(sw / ratio)
                else:
                    final_h = sh
                    final_w = int(sh * ratio)
                left = int((sw - final_w) / 2)
                top  = int((sh - final_h) / 2)

            slide.shapes.add_picture(img_byte_arr, left, top, width=final_w, height=final_h)

        # 「拡張子を除いたファイル名」を左上 (0,0) にテキストボックス表示
        name_no_ext = os.path.splitext(image_name)[0]
        textbox = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.RECTANGLE,
            left=0,
            top=0,
            width=int(prs.slide_width * 0.3),
            height=int(prs.slide_height * 0.06),
        )
        # ボックス背景：オレンジ
        textbox.fill.solid()
        textbox.fill.fore_color.rgb = RGBColor(255, 102, 0)
        # 枠線：赤
        textbox.line.color.rgb = RGBColor(255, 0, 0)
        # 影を無効化
        textbox.shadow.inherit = False

        # テキスト設定
        text_frame = textbox.text_frame
        p = text_frame.paragraphs[0]
        p.text = name_no_ext
        font = p.font
        font.name = "游ゴシック"
        font.size = Pt(18)
        font.bold = True
        font.color.rgb = RGBColor(255, 255, 255)

    save_path = os.path.join(base_dir, "Slides_from_Images.pptx")
    prs.save(save_path)
    messagebox.showinfo("完了", f"PPTX(Images→A3実寸)作成が完了しました\n\n保存先: {save_path}")
    webbrowser.open(save_path)

# ===============================================
# 機能 3：PDF→PPTX 直接変換 (中間PNGを保持せずに実寸で配置)
# ===============================================

def create_pptx_direct_from_pdfs():
    """
    Inputフォルダ内のPDFをすべて読み込み、各ページを
    A3横スライドに「元のサイズ(実寸)」で中央配置し
    拡張子を除いたファイル名テキストを左上に表示
    1つのPPTXにまとめる。中間PNGファイルは生成せず扱う。
    プログレスバーでページごとに進捗を表示。
    """
    base_dir = get_app_path()
    input_dir = os.path.join(base_dir, "Input")
    pdf_files = [f for f in os.listdir(input_dir) if f.lower().endswith(".pdf")]
    if not pdf_files:
        messagebox.showwarning("警告", "InputフォルダにPDFが見つかりません")
        return

    # 全ページ数をカウント (プログレスバー用)
    total_pages = 0
    for fname in pdf_files:
        try:
            doc = fitz.open(os.path.join(input_dir, fname))
            total_pages += len(doc)
        except:
            continue

    # スケール取得（画像解像度向上のための倍率。ただし配置は実寸）
    try:
        scale = float(entry_scale.get())
    except ValueError:
        messagebox.showerror("エラー", "スケールは数値で指定してください")
        return

    # プレゼンテーション作成＆A3横設定
    prs = Presentation()
    prs.slide_width  = mm_to_emu(420)
    prs.slide_height = mm_to_emu(297)

    # プログレスバー初期化
    progress["value"] = 0
    progress["maximum"] = total_pages
    root.update_idletasks()

    for file_name in pdf_files:
        file_path = os.path.join(input_dir, file_name)
        try:
            doc = fitz.open(file_path)
            base_name = os.path.splitext(file_name)[0]
            for page_number in range(len(doc)):
                page = doc[page_number]

                # PDFページの実寸サイズ（ポイント単位→EMU）
                rect = page.rect
                width_emu  = points_to_emu(rect.width)
                height_emu = points_to_emu(rect.height)

                # 画像を BytesIO 経由で取得 (Zoom=scale をかけて解像度確保)
                zoom_mtx = fitz.Matrix(scale, scale)
                pix = page.get_pixmap(matrix=zoom_mtx)
                png_bytes = pix.tobytes("png")
                image_stream = BytesIO(png_bytes)

                # 図を配置する位置は、スライド中央にするためのオフセット計算
                sw = prs.slide_width
                sh = prs.slide_height
                final_w = width_emu
                final_h = height_emu
                # スライドより大きければフィット縮小
                if final_w > sw or final_h > sh:
                    ratio = final_w / final_h
                    slide_ratio = sw / sh
                    if ratio > slide_ratio:
                        final_w = sw
                        final_h = int(sw / ratio)
                    else:
                        final_h = sh
                        final_w = int(sh * ratio)
                left = int((sw - final_w) / 2)
                top  = int((sh - final_h) / 2)

                # 空白スライド追加
                slide = prs.slides.add_slide(prs.slide_layouts[6])
                slide.shapes.add_picture(image_stream, left, top, width=final_w, height=final_h)

                # 「拡張子を除いたファイル名_ページ番号」を左上に表示
                display_name = f"{base_name}_{page_number + 1}"
                textbox = slide.shapes.add_shape(
                    MSO_AUTO_SHAPE_TYPE.RECTANGLE,
                    left=0,
                    top=0,
                    width=int(prs.slide_width * 0.3),
                    height=int(prs.slide_height * 0.06),
                )
                textbox.fill.solid()
                textbox.fill.fore_color.rgb = RGBColor(255, 102, 0)
                textbox.line.color.rgb = RGBColor(255, 0, 0)
                # 影を無効化
                textbox.shadow.inherit = False

                tf = textbox.text_frame
                p = tf.paragraphs[0]
                p.text = display_name
                font = p.font
                font.name = "游ゴシック"
                font.size = Pt(18)
                font.bold = True
                font.color.rgb = RGBColor(255, 255, 255)

                # プログレス更新
                progress["value"] += 1
                root.update_idletasks()

        except Exception as e:
            print(f"Error processing {file_name}: {e}")

    save_path = os.path.join(base_dir, "Slides_from_PDF_direct.pptx")
    prs.save(save_path)
    messagebox.showinfo("完了", f"PPTX(PDF→A3実寸)作成が完了しました\n\n保存先: {save_path}")
    webbrowser.open(save_path)
    progress["value"] = 0

# ===============================================
# 機能 4：Input/Output 初期化
# ===============================================

def clear_folders():
    """
    Input, Output フォルダを完全に削除し、空のフォルダを再作成する。
    """
    base_dir = get_app_path()
    for name in ("Input", "Output"):
        target = os.path.join(base_dir, name)
        if os.path.exists(target):
            try:
                shutil.rmtree(target)
            except Exception as e:
                messagebox.showerror("削除エラー", f"{target}: {e}")
        os.makedirs(target)
    messagebox.showinfo("初期化", "Input / Output フォルダを初期化しました")

# ===============================================
# メイン：GUI の構築
# ===============================================

# 必要なフォルダを作成
ensure_dirs()

root = tk.Tk()
root.title("PDFToolkit (PDF→PNG→PPTX／PDF→PPTX実寸)")

# スケール倍率設定
tk.Label(root, text="スケール倍率 (解像度向上 例: 1.5):").grid(row=0, column=0, sticky="e", padx=10, pady=10)
entry_scale = tk.Entry(root, width=10)
entry_scale.insert(0, "1.5")
entry_scale.grid(row=0, column=1, padx=10)

# ボタン群
btn1 = tk.Button(root, text="📄 画像変換 (PDF → PNG)", command=convert_pdfs_to_images, width=30)
btn1.grid(row=1, column=0, columnspan=2, pady=5)

btn2 = tk.Button(root, text="📊 PPTX作成 (PNG→A3実寸)", command=create_pptx_from_images, width=30)
btn2.grid(row=2, column=0, columnspan=2, pady=5)

btn3 = tk.Button(root, text="📈 PPTX作成 (PDF→A3実寸 直接)", command=create_pptx_direct_from_pdfs, width=30)
btn3.grid(row=3, column=0, columnspan=2, pady=5)

btn4 = tk.Button(root, text="🧹 Input/Output 初期化", command=clear_folders, width=30)
btn4.grid(row=4, column=0, columnspan=2, pady=5)

# プログレスバー
progress = ttk.Progressbar(root, orient="horizontal", length=400, mode="determinate")
progress.grid(row=5, column=0, columnspan=2, pady=10)

root.mainloop()

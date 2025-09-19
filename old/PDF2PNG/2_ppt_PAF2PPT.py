import fitz  # PyMuPDF
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.slide import SlideLayout

def pdfs_to_single_ppt(input_folder, scale=2.0, output_pptx_path="output_presentation.pptx"):
    """
    複数のPDFを1つのPPTXにまとめ、画像を中央配置かつサイズ調整する（空白スライドに画像のみを貼り付け）。
    
    Parameters:
        input_folder (str): PDFファイルが格納されているフォルダ
        scale (float): 画像スケール。1.0で元サイズ、2.0で2倍。
        output_pptx_path (str): 出力するPPTXファイルのパス
    """
    prs = Presentation()  # 新しいPPTXを作成
    slide_width, slide_height = prs.slide_width, prs.slide_height  # スライドのサイズ

    for file_name in os.listdir(input_folder):
        if file_name.lower().endswith('.pdf'):
            file_path = os.path.join(input_folder, file_name)
            try:
                doc = fitz.open(file_path)  # PDFを開く
                base_name = os.path.splitext(file_name)[0]

                for page_number in range(len(doc)):
                    page = doc[page_number]

                    # サイズを調整するためのスケール設定
                    zoom_matrix = fitz.Matrix(scale, scale)
                    pix = page.get_pixmap(matrix=zoom_matrix)
                    
                    # 画像サイズを取得
                    img_width = Inches(pix.width / 96)  # 96 DPI 基準でインチに変換
                    img_height = Inches(pix.height / 96)

                    # 画像のスケーリング比率を計算してスライド内に収める
                    scale_ratio = min(slide_width / img_width, slide_height / img_height, 1.0)
                    img_width *= scale_ratio
                    img_height *= scale_ratio
                    
                    # 画像を中央に配置
                    left = (slide_width - img_width) / 2
                    top = (slide_height - img_height) / 2

                    # 一時画像ファイルに保存
                    temp_image_path = os.path.join(input_folder, f"{base_name}_{page_number + 1}.png")
                    pix.save(temp_image_path)
                    
                    # 空白スライドを追加（完全に空白のスライド）
                    slide_layout = prs.slide_layouts[6]  # 空白レイアウト (6番目のレイアウトが空白)
                    slide = prs.slides.add_slide(slide_layout)
                    
                    # スライドに画像を貼り付け
                    slide.shapes.add_picture(temp_image_path, left, top, width=img_width, height=img_height)
                    
                    # 一時画像ファイルを削除
                    os.remove(temp_image_path)

            except Exception as e:
                print(f"Error processing {file_name}: {e}")
    
    # PPTXを保存
    prs.save(output_pptx_path)
    print(f"Saved PPTX: {output_pptx_path}")

# 実行
input_folder = './Input'
scale_factor = 1.5  # 画像のスケールを1.5倍に設定
output_pptx_path = './Output/combined_presentation.pptx'
pdfs_to_single_ppt(input_folder, scale=scale_factor, output_pptx_path=output_pptx_path)
